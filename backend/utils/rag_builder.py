import os
import json
import re
from io import BytesIO

import pandas as pd
import pytesseract
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from langchain_core.documents import Document
from PIL import Image, ImageFilter, ImageOps, ImageStat
from pptx import Presentation
from pytesseract import TesseractNotFoundError
from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

FAISS_INDEX_DIR = os.getenv("FAISS_INDEX_DIR", "./data/faiss_indexes")
EMBEDDING_MODEL = "all-MiniLM-L6-v2"
SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".txt",
    ".md",
    ".docx",
    ".csv",
    ".xlsx",
    ".json",
    ".html",
    ".htm",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".tiff",
}

_IMAGE_MIME_TYPES = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
    ".tiff": "image/tiff",
}

# Minimum edge-variance score for an image to be considered clear enough for OCR.
# Increase this value to require higher image sharpness.
_MIN_CLARITY_SCORE = float(os.getenv("MIN_IMAGE_CLARITY", "30.0"))


def _get_embeddings() -> HuggingFaceEmbeddings:
    """Return a HuggingFace embedding model instance (cached locally)."""
    return HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)


def get_index_path(user_id: int, document_id: int) -> str:
    """Return the filesystem path for a user's document FAISS index."""
    return os.path.join(FAISS_INDEX_DIR, str(user_id), str(document_id))


def _extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(file_bytes))
    pages = [(page.extract_text() or "") for page in reader.pages]
    return "\n\n".join(pages).strip()


def _extract_text_from_text_like(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore").strip()


def _extract_text_from_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text).strip()


def _extract_text_from_csv(file_bytes: bytes) -> str:
    df = pd.read_csv(BytesIO(file_bytes), dtype=str, na_filter=False)
    return df.to_csv(index=False)


def _extract_text_from_xlsx(file_bytes: bytes) -> str:
    sheets = pd.read_excel(BytesIO(file_bytes), sheet_name=None, dtype=str)
    parts: list[str] = []
    for name, df in sheets.items():
        parts.append(f"# Sheet: {name}\n{df.fillna('').to_csv(index=False)}")
    return "\n\n".join(parts).strip()


def _extract_text_from_json(file_bytes: bytes) -> str:
    parsed = json.loads(file_bytes.decode("utf-8", errors="ignore"))
    return json.dumps(parsed, indent=2, ensure_ascii=True)


def _extract_text_from_html(file_bytes: bytes) -> str:
    soup = BeautifulSoup(file_bytes, "html.parser")
    return soup.get_text("\n", strip=True)


def _extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                blocks.append(shape.text)
        if blocks:
            slides.append(f"# Slide {i}\n" + "\n".join(blocks))
    return "\n\n".join(slides).strip()


def _is_low_value_ocr_text(text: str) -> bool:
    normalized = re.sub(r"\s+", " ", text or "").strip().lower()
    if not normalized:
        return True
    alnum_len = len(re.sub(r"[^a-z0-9]+", "", normalized))
    return alnum_len < 8


def _run_local_ocr(image: Image.Image) -> str:
    variants = [
        image,
        image.convert("L"),
        ImageOps.autocontrast(image.convert("L")),
    ]
    # Binary threshold often improves scanned docs.
    gray = ImageOps.autocontrast(image.convert("L"))
    bw = gray.point(lambda p: 255 if p > 145 else 0)
    variants.append(bw)

    attempts: list[str] = []
    for variant in variants:
        text = pytesseract.image_to_string(variant, config="--oem 3 --psm 6").strip()
        if text:
            attempts.append(text)

    if not attempts:
        return ""

    # Prefer richest output; it is usually closest to full-page extraction.
    attempts.sort(key=len, reverse=True)
    return attempts[0]


def _check_image_clarity(image: Image.Image) -> float:
    """Return a sharpness score (variance of edge pixels). Higher = sharper."""
    gray = image.convert("L")
    edges = gray.filter(ImageFilter.FIND_EDGES)
    return ImageStat.Stat(edges).var[0]


def _extract_text_from_image(file_bytes: bytes, ext: str) -> str:
    try:
        image = Image.open(BytesIO(file_bytes))
        image.load()
    except Exception as exc:
        raise ValueError(
            "Could not open this image. Please upload a valid image file (PNG, JPG, etc.)."
        ) from exc

    # Reject blurry or low-quality images before attempting OCR.
    clarity = _check_image_clarity(image)
    if clarity < _MIN_CLARITY_SCORE:
        raise ValueError(
            f"This image is too blurry or low quality (clarity score: {clarity:.1f}, "
            f"minimum required: {_MIN_CLARITY_SCORE:.1f}). "
            "Please upload a clearer, higher-resolution scan or photo with readable text."
        )

    # Run OCR — only text-bearing images are supported.
    try:
        ocr_text = _run_local_ocr(image)
    except TesseractNotFoundError as exc:
        raise ValueError(
            "OCR is not available on this server. "
            "Contact the administrator to install Tesseract."
        ) from exc

    if _is_low_value_ocr_text(ocr_text):
        raise ValueError(
            "No readable text was found in this image. "
            "Only images that contain clear, printed or typed text are supported. "
            "Please upload a document scan or a photo of text-based content."
        )

    return ocr_text.strip()


def _extract_text(file_bytes: bytes, ext: str) -> str:
    if ext == ".pdf":
        return _extract_text_from_pdf(file_bytes)
    if ext in {".txt", ".md"}:
        return _extract_text_from_text_like(file_bytes)
    if ext == ".docx":
        return _extract_text_from_docx(file_bytes)
    if ext == ".csv":
        return _extract_text_from_csv(file_bytes)
    if ext == ".xlsx":
        return _extract_text_from_xlsx(file_bytes)
    if ext == ".json":
        return _extract_text_from_json(file_bytes)
    if ext in {".html", ".htm"}:
        return _extract_text_from_html(file_bytes)
    if ext == ".pptx":
        return _extract_text_from_pptx(file_bytes)
    if ext in _IMAGE_MIME_TYPES:
        return _extract_text_from_image(file_bytes, ext)
    raise ValueError(f"Unsupported extension: {ext}")


def build_faiss_index(file_bytes: bytes, ext: str, user_id: int, document_id: int) -> str:
    """
    Accept raw file bytes, split into overlapping chunks, embed with a local
    HuggingFace model, build a FAISS vector store, and persist it to disk.

    Supports: PDF/TXT/MD/DOCX/CSV/XLSX/JSON/HTML/PPTX and common image formats.
    Returns: absolute path to the saved FAISS index directory.
    """
    ext = ext.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type '{ext}'. Supported types: {supported}")

    extracted_text = _extract_text(file_bytes, ext)
    if not extracted_text:
        raise ValueError("No extractable text found in the uploaded file.")

    documents = [Document(page_content=extracted_text, metadata={"source_ext": ext})]

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
    )
    chunks = splitter.split_documents(documents)

    if not chunks:
        raise ValueError("No extractable text chunks were generated from the uploaded file.")

    embeddings = _get_embeddings()
    vector_store = FAISS.from_documents(chunks, embeddings)

    index_path = get_index_path(user_id, document_id)
    os.makedirs(index_path, exist_ok=True)
    vector_store.save_local(index_path)

    return index_path


def load_faiss_index(user_id: int, document_id: int) -> FAISS:
    """
    Load a persisted FAISS index from disk for similarity search / RAG querying.
    """
    index_path = get_index_path(user_id, document_id)
    if not os.path.exists(index_path):
        raise FileNotFoundError(f"FAISS index not found at: {index_path}")
    embeddings = _get_embeddings()
    return FAISS.load_local(index_path, embeddings, allow_dangerous_deserialization=True)
